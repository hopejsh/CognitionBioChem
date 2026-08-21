#!/usr/bin/env python3
"""Render a generated HTML deck as an editable PowerPoint file.

    ./.venv/bin/python platform/build_pptx.py                       # the workbench deck
    ./.venv/bin/python platform/build_pptx.py <deck.html> [out.pptx] # any deck this repo makes

Decks in this repository are authored once as HTML and re-rendered; this module is the
PowerPoint renderer, and it takes whichever deck it is pointed at.

This does NOT re-author the slides. It parses docs/CognitionBioChem_Deck.html -- the deck
build_deck.py just produced -- and re-renders it as native PowerPoint shapes: text frames,
tables and pictures, all editable. There is therefore exactly one place the deck's prose
lives, and the .pptx cannot drift from the .html or the .pdf. If build_deck.py changes the
markup in a way this parser does not recognise, it raises rather than silently dropping a
slide.

Fonts. The deck is set in IBM Plex Sans Condensed / Newsreader / IBM Plex Mono, which are web
fonts and are not installed locally. PowerPoint would substitute them silently and badly, so
the .pptx names the closest faces that macOS actually ships:

    display  IBM Plex Sans Condensed  ->  Avenir Next Condensed  (bold)
    body     Newsreader               ->  Iowan Old Style
    data     IBM Plex Mono            ->  Menlo

On a machine without those -- Windows, or a Mac with a trimmed font set -- PowerPoint
substitutes again. Change FONT_DISPLAY / FONT_BODY / FONT_MONO below and re-run if you would
rather it named different faces.
"""

from __future__ import annotations

import base64
import io
import re
import sys

from html.parser import HTMLParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "platform"))
from cbc.provenance import git_sha  # noqa: E402

DEFAULT_SRC = REPO / "docs" / "CognitionBioChem_Deck.html"

FONT_DISPLAY = "Avenir Next Condensed"
FONT_BODY = "Iowan Old Style"
FONT_MONO = "Menlo"

#: None of the three faces above carries a single Hangul glyph, so a Korean deck rendered with
#: them names a font that cannot draw its own text and PowerPoint substitutes silently --
#: differently on every machine. A deck whose markup declares lang="ko" gets faces that cover
#: it. Menlo stays for the mono column, which carries Latin identifiers and numerals only.
FONT_KO_DISPLAY = "Apple SD Gothic Neo"
FONT_KO_BODY = "Apple SD Gothic Neo"
FONT_KO_MONO = "Menlo"

# The deck is laid out in a 1120x630 CSS-pixel frame. One deck pixel is 1/84 inch, so every
# position and size below can be written in the same units the stylesheet uses.
DECK_W, DECK_H = 1120, 630
PX = Inches(1 / 84)
PT_PER_PX = 72 / 84          # 0.857 pt per deck pixel


def px(v: float) -> Emu:
    return Emu(int(round(v * PX)))


def pt(v: float) -> Pt:
    return Pt(round(v * PT_PER_PX, 1))


# Light palette only: a .pptx has one appearance, and a deck that is printed, projected and
# emailed should be the light one.
INK = RGBColor(0x17, 0x1A, 0x19)
MUTED = RGBColor(0x6C, 0x72, 0x6D)
RULE = RGBColor(0xD8, 0xDB, 0xD4)
PANEL = RGBColor(0xFB, 0xFB, 0xF8)
ACCENT = RGBColor(0x2F, 0x6F, 0x62)
AMBER = RGBColor(0xB0, 0x76, 0x1C)
SLATE = RGBColor(0x5A, 0x6B, 0x7C)
RAIL = {"premise": RGBColor(0x8A, 0x8F, 0x88), "method": ACCENT,
        "result": AMBER, "limit": SLATE}
CHIP = {"c": ACCENT, "f": SLATE, "n": MUTED}


# ============================================================== parsing the deck ======= #
#: Inline styles the deck's markup can put on a run of text.
_TAG_STYLE = {"b": "bold", "strong": "bold", "i": "italic", "em": "italic", "sub": "sub"}
_CLASS_STYLE = {"mono": "mono", "hl": "amber", "hl-a": "accent", "strike": "strike",
                "sym": None, "k": "key", "kind": "kind", "stamp": "stamp", "no": "no",
                "chip": "chip", "c": "chip-c", "f": "chip-f", "n": "chip-n",
                "thesis": None, "meta": None, "ids": None}


class _Inline(HTMLParser):
    """Flatten a fragment of the deck's markup into styled runs.

    Runs are (text, frozenset_of_styles). Everything the deck uses inline is covered; an
    unrecognised tag raises, because a silently dropped <b> is a defect that survives review.
    """

    ALLOWED = {"b", "strong", "i", "em", "sub", "span", "br", "p", "div", "ul", "li",
               "table", "thead", "tbody", "tr", "th", "td", "img", "figcaption", "a"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.runs: list[tuple[str, frozenset]] = []
        self._stack: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag not in self.ALLOWED:
            raise ValueError(f"inline parser met an unhandled <{tag}>")
        if tag == "br":
            self.runs.append(("\n", frozenset()))
            return
        st = _TAG_STYLE.get(tag)
        if tag == "span":
            cls = dict(attrs).get("class", "").split()
            st = next((_CLASS_STYLE[c] for c in cls if _CLASS_STYLE.get(c)), None)
            if "chip" in cls:
                st = "chip-" + next((c for c in cls if c in "cfn"), "n")
        self._stack.append(st)

    def handle_endtag(self, tag):
        if tag == "br":
            return
        if self._stack:
            self._stack.pop()

    def handle_data(self, data):
        text = re.sub(r"\s+", " ", data)
        if text:
            self.runs.append((text, frozenset(s for s in self._stack if s)))

    @staticmethod
    def unescape(text: str) -> str:
        import html as _h
        return _h.unescape(text)

    @classmethod
    def runs_of(cls, html: str) -> list[tuple[str, frozenset]]:
        p = cls()
        p.feed(html)
        p.close()
        out: list[tuple[str, frozenset]] = []
        for text, style in p.runs:                       # merge adjacent same-styled runs
            if out and out[-1][1] == style:
                out[-1] = (out[-1][0] + text, style)
            else:
                out.append((text, style))
        if out:
            out[0] = (out[0][0].lstrip(), out[0][1])
            out[-1] = (out[-1][0].rstrip(), out[-1][1])
        return [r for r in out if r[0]]


def _tag_span(html: str, start: int) -> tuple[str, int]:
    """Return the inner HTML of the element opening at `start`, and the index after it.

    A regex cannot do this: the deck nests divs inside divs. This walks the tag stream and
    counts depth, which is enough for markup this generator itself emits.
    """
    m = re.compile(r"<(\w+)").match(html, start)
    if not m:
        raise ValueError(f"no element at offset {start}: {html[start:start + 60]!r}")
    name = m.group(1)
    open_end = html.index(">", start) + 1
    if html[open_end - 2] == "/":
        return "", open_end
    depth, i = 1, open_end
    pat = re.compile(rf"</?{name}\b", re.I)
    while depth:
        mm = pat.search(html, i)
        if not mm:
            raise ValueError(f"unbalanced <{name}> from offset {start}")
        if html[mm.start() + 1] == "/":
            depth -= 1
            i = html.index(">", mm.start()) + 1
            if depth == 0:
                return html[open_end:mm.start()], i
        else:
            depth += 1
            i = html.index(">", mm.start()) + 1
    raise AssertionError


def _attrs(open_tag: str) -> dict:
    return dict(re.findall(r'(\w[\w-]*)="([^"]*)"', open_tag))


def parse_slides(html: str) -> list[dict]:
    """Split the deck into slides and each slide into typed blocks."""
    slides = []
    for m in re.finditer(r'<section class="slide([^"]*)"([^>]*)>', html):
        inner, _ = _tag_span(html, m.start())
        attrs = _attrs(m.group(0))
        slides.append({
            "kind": attrs.get("data-kind", "premise"),
            "title_slide": "title" in m.group(1).split(),
            "bleed": attrs.get("data-fig") == "bleed",
            "blocks": _parse_slide(inner),
        })
    if not slides:
        raise ValueError("no <section class=\"slide\"> found -- has build_deck.py changed?")
    return slides


def _parse_slide(inner: str) -> list[dict]:
    blocks: list[dict] = []
    i = 0
    while i < len(inner):
        m = re.compile(r"<(\w+)([^>]*)>").search(inner, i)
        if not m:
            break
        tag, attrs = m.group(1), _attrs(m.group(0))
        cls = attrs.get("class", "").split()
        body, after = _tag_span(inner, m.start())
        if tag in ("header", "footer") or "eyebrow" in cls or "meta" in cls:
            blocks.append({"t": tag if tag in ("header", "footer") else cls[0],
                           "html": body})
        elif tag in ("h1", "h2"):
            blocks.append({"t": "h", "level": int(tag[1]), "html": body})
        elif tag == "div" and "body" in cls:
            blocks.extend(_parse_body(body))
        else:
            blocks.extend(_parse_body(inner[m.start():after]))
        i = after
    return blocks


def _parse_body(html: str) -> list[dict]:
    """Parse the block-level content of a slide body (or of one column of it)."""
    out: list[dict] = []
    i = 0
    while i < len(html):
        m = re.compile(r"<(\w+)([^>]*)>").search(html, i)
        if not m:
            break
        tag, attrs = m.group(1), _attrs(m.group(0))
        cls = attrs.get("class", "").split()
        body, after = _tag_span(html, m.start())

        if "cols" in cls:
            ratio = (7, 5) if "c-7-5" in cls else (1, 1)
            cols, j = [], 0
            while j < len(body):
                cm = re.compile(r"<(\w+)").search(body, j)
                if not cm:
                    break
                cb, ca = _tag_span(body, cm.start())
                cols.append(_parse_body(body[cm.start():ca]))
                j = ca
            out.append({"t": "cols", "ratio": ratio, "cols": cols})
        elif "stats" in cls:
            stats = []
            for sm in re.finditer(r'<div class="stat">', body):
                sb, _ = _tag_span(body, sm.start())
                b = re.search(r"<b([^>]*)>", sb)
                bb, ba = _tag_span(sb, b.start())
                lb = re.search(r"<span([^>]*)>", sb[ba:])
                lbody, _ = _tag_span(sb[ba:], lb.start())
                stats.append({"value": _Inline.runs_of(bb),
                              "label": _Inline.runs_of(lbody),
                              "struck": "strike" in _attrs(b.group(0)).get("class", "")})
            out.append({"t": "stats", "items": stats})
        elif "figwrap" in cls:
            img = re.search(r'<img src="data:image/(\w+);base64,([^"]+)"', body)
            cap = re.search(r"<figcaption[^>]*>", body)
            capbody = _tag_span(body, cap.start())[0] if cap else ""
            out.append({"t": "fig",
                        "png": base64.b64decode(img.group(2)) if img else None,
                        "caption": _Inline.runs_of(capbody)})
        elif tag == "ul" and "led" in cls:
            items = []
            for lm in re.finditer(r"<li[^>]*>", body):
                lb, _ = _tag_span(body, lm.start())
                spans = []
                j2 = 0
                while j2 < len(lb):
                    sm = re.compile(r"<span([^>]*)>").search(lb, j2)
                    if not sm:
                        break
                    sb, sa = _tag_span(lb, sm.start())
                    spans.append((_attrs(sm.group(0)).get("class", ""), sb))
                    j2 = sa
                items.append({"key": _Inline.runs_of(spans[0][1]) if spans else [],
                              "text": _Inline.runs_of(spans[1][1]) if len(spans) > 1 else []})
            out.append({"t": "led", "items": items})
        elif tag == "table" or "tablewrap" in cls:
            tb = body if tag == "table" else _tag_span(body, body.index("<table"))[0]
            # `<th[^>]*>` also matches the opening `<thead ...>`, which swallowed the
            # whole header row as a phantom first column. Require a tag boundary.
            head = [_Inline.runs_of(_tag_span(tb, hm.start())[0])
                    for hm in re.finditer(r"<th(?=[\s>])", tb)]
            rows = []
            for rm in re.finditer(r"<tr[^>]*>", tb):
                rb, _ = _tag_span(tb, rm.start())
                if "<th" in rb:
                    continue
                cells, j3 = [], 0
                while j3 < len(rb):
                    cm = re.compile(r"<td(?=[\s>])([^>]*)>").search(rb, j3)
                    if not cm:
                        break
                    cb, ca = _tag_span(rb, cm.start())
                    cells.append({"runs": _Inline.runs_of(cb),
                                  "attrs": _attrs(cm.group(0))})
                    j3 = ca
                if cells:
                    rows.append(cells)
            out.append({"t": "table", "head": head, "rows": rows})
        elif tag in ("h1", "h2"):
            # A title slide wraps its headline in a plain <div>, so the heading arrives here
            # rather than at slide level. Dropping it cost the deck its two title slides.
            out.append({"t": "h", "level": int(tag[1]), "html": body})
        elif tag == "p":
            out.append({"t": "pull" if "pull" in cls else "p",
                        "html": body, "runs": _Inline.runs_of(body)})
        elif tag == "div":
            out.extend(_parse_body(body))
        i = after
    return out


# ============================================================== rendering to pptx ====== #
#: Average advance width as a fraction of the point size, per face. Used only to estimate
#: how many lines a run of text will take, so blocks can be flowed without overlapping.
#: PowerPoint does the real line breaking; these numbers decide where the NEXT block starts.
_ADVANCE = {FONT_DISPLAY: 0.46, FONT_BODY: 0.478, FONT_MONO: 0.615}
LINE = 1.35


def est_height(runs, width_px: float, size_px: float, face: str | None = None,
               line: float = LINE) -> float:
    """Height in deck pixels that `runs` will occupy in a box `width_px` wide."""
    face = face or FONT_BODY
    if not runs:
        return 0.0
    per_char = _ADVANCE.get(face, 0.46) * size_px
    text = "".join(t for t, _ in runs)
    lines = 0
    for hard in text.split("\n"):
        lines += max(1, -(-len(hard) * per_char // width_px))
    return lines * size_px * line


def _spacing(run, hundredths: int) -> None:
    """Letter-spacing, which python-pptx does not expose."""
    run.font._rPr.set("spc", str(int(hundredths)))


def box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def put(para, runs, *, size_px, face=None, color=INK, bold=False, italic=False,
        line=LINE, space_after_px=0, align=PP_ALIGN.LEFT, spc=None, upper=False,
        strike=False):
    """Write styled runs into a paragraph, honouring the deck's inline styles."""
    face = face or FONT_BODY
    para.alignment = align
    para.line_spacing = line
    para.space_after = px(space_after_px)
    para.space_before = 0
    for text, style in runs or [("", frozenset())]:
        r = para.add_run()
        r.text = text.upper() if upper else text
        f = r.font
        f.size = pt(size_px)
        f.name = FONT_MONO if "mono" in style or "key" in style else face
        f.bold = bold or "bold" in style
        f.italic = italic or "italic" in style
        f.color.rgb = color
        if "amber" in style:
            f.color.rgb, f.bold = AMBER, True
        if "accent" in style or "key" in style:
            f.color.rgb = ACCENT
        if strike or "strike" in style:
            f.color.rgb = MUTED
            f._rPr.set("strike", "sngStrike")
        for chip, col in CHIP.items():
            if f"chip-{chip}" in style:
                f.color.rgb, f.name, f.size = col, FONT_MONO, pt(size_px * 0.88)
        if "sub" in style:
            f._rPr.set("baseline", "-25000")
        if spc:
            _spacing(r, spc)
        # A mono run inside body text reads as data, so it keeps the body colour but never
        # the body face; the deck makes the same distinction with .mono.
        if "mono" in style and not (style & {"amber", "accent", "key"}):
            f.color.rgb = color
    return para


def rule(slide, x, y, w, color=RULE, h_px=0.8):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h_px))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


# --- geometry, in deck pixels ------------------------------------------------------------
L, R, TOP, BOT = 50.0, 40.0, 34.0, 24.0
CW = DECK_W - L - R                     # content width
FOOT_Y = DECK_H - BOT - 14.0
GAP = 26.0                              # the deck's column gap

SIZE = {"eyebrow": 9.3, "h2": 34.0, "h1": 57.0, "sub": 14.3, "p": 13.1, "pull": 17.9,
        "stat": 30.0, "statlabel": 8.6, "foot": 8.8, "table": 11.0, "cap": 8.8,
        "thesis": 18.0, "meta": 9.3, "key": 8.6}


def _runs(html: str):
    return _Inline.runs_of(html)


def _eyebrow(slide, html: str, y: float) -> float:
    """The typed label, the section gloss, and the registered plan hash on the right."""
    parts = re.findall(r'<span([^>]*)>(.*?)</span>', html, re.S)
    kind = next((t for a, t in parts if "kind" in a), "")
    stamp = next((t for a, t in parts if "stamp" in a), "")
    gloss = [t for a, t in parts if "kind" not in a and "stamp" not in a]
    _, tf = box(slide, L, y, CW * 0.72, 14)
    p = tf.paragraphs[0]
    put(p, _runs(kind), size_px=SIZE["eyebrow"], face=FONT_MONO, color=INK, bold=True,
        line=1.0, spc=120, upper=True)
    for g in gloss:
        put(p, [("   ", frozenset())] + _runs(g), size_px=SIZE["eyebrow"], face=FONT_MONO,
            color=MUTED, line=1.0, spc=120, upper=True)
    if stamp:
        _, tf2 = box(slide, L + CW * 0.72, y, CW * 0.28, 14)
        put(tf2.paragraphs[0], _runs(stamp), size_px=SIZE["eyebrow"], face=FONT_MONO,
            color=MUTED, line=1.0, spc=60, align=PP_ALIGN.RIGHT)
    return y + 20


def _footer(slide, html: str) -> None:
    parts = re.findall(r'<span([^>]*)>(.*?)</span>', html, re.S)
    left = " ".join(t for a, t in parts if "no" not in a)
    num = next((t for a, t in parts if 'class="no"' in a), "")
    rule(slide, L, FOOT_Y - 9, CW)
    _, tf = box(slide, L, FOOT_Y, CW * 0.9, 14)
    put(tf.paragraphs[0], _runs(left), size_px=SIZE["foot"], face=FONT_MONO, color=MUTED,
        line=1.0, spc=45)
    _, tf2 = box(slide, L + CW * 0.9, FOOT_Y, CW * 0.1, 14)
    put(tf2.paragraphs[0], _runs(num), size_px=SIZE["foot"], face=FONT_MONO, color=MUTED,
        line=1.0, spc=45, align=PP_ALIGN.RIGHT)


def _header(slide, blocks, y: float) -> float:
    """Eyebrow, headline and standfirst. Returns the y the body may start at."""
    html = next((b["html"] for b in blocks if b["t"] == "header"), None)
    if html is None:
        return y
    eb = re.search(r'<div class="eyebrow">', html)
    if eb:
        y = _eyebrow(slide, _tag_span(html, eb.start())[0], y)
    hm = re.search(r"<h2[^>]*>", html)
    if hm:
        runs = _runs(_tag_span(html, hm.start())[0])
        h = est_height(runs, CW, SIZE["h2"], FONT_DISPLAY, 1.06)
        _, tf = box(slide, L, y, CW, h + 4)
        put(tf.paragraphs[0], runs, size_px=SIZE["h2"], face=FONT_DISPLAY, bold=True,
            line=1.06)
        y += h + 6
    sm = re.search(r'<p class="sub">', html)
    if sm:
        runs = _runs(_tag_span(html, sm.start())[0])
        h = est_height(runs, CW * 0.62, SIZE["sub"], FONT_BODY, 1.35)
        _, tf = box(slide, L, y, CW * 0.62, h + 3)
        put(tf.paragraphs[0], runs, size_px=SIZE["sub"], color=MUTED, line=1.35)
        y += h + 6
    return y + 16


def _blocks(slide, blocks, x: float, y: float, w: float, avail: float) -> float:
    """Lay out a column of body blocks from (x, y) downwards. Returns the new y."""
    i = 0
    while i < len(blocks):
        b = blocks[i]
        if b["t"] == "p" and b["runs"]:
            group = []
            while i < len(blocks) and blocks[i]["t"] == "p":
                if blocks[i]["runs"]:
                    group.append(blocks[i]["runs"])
                i += 1
            h = sum(est_height(g, w, SIZE["p"], FONT_BODY, LINE) for g in group)
            h += 9 * (len(group) - 1)
            _, tf = box(slide, x, y, w, h + 6)
            for k, g in enumerate(group):
                para = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                put(para, g, size_px=SIZE["p"], line=LINE,
                    space_after_px=9 if k < len(group) - 1 else 0)
            y += h + 13
            continue
        i += 1
        if b["t"] == "cols":
            a, c = b["ratio"]
            wa = (w - GAP) * a / (a + c)
            wc = (w - GAP) * c / (a + c)
            ya = _blocks(slide, b["cols"][0], x, y, wa, avail)
            yc = _blocks(slide, b["cols"][1], x + wa + GAP, y, wc, avail) \
                if len(b["cols"]) > 1 else y
            y = max(ya, yc)

        elif b["t"] == "pull":
            runs = b["runs"]
            if not runs:
                continue
            h = est_height(runs, w - 14, SIZE["pull"], FONT_DISPLAY, 1.24)
            rule(slide, x, y, 2.6, RAIL["method"], h)
            _, tf = box(slide, x + 14, y, w - 14, h + 4)
            put(tf.paragraphs[0], runs, size_px=SIZE["pull"], face=FONT_DISPLAY, bold=True,
                line=1.24)
            y += h + 18

        elif b["t"] == "stats":
            items = b["items"]
            ncol = 2 if w < DECK_W * 0.55 else min(len(items), 6)
            cw = (w - GAP * (ncol - 1)) / ncol
            # Measure every label first. Placing while the row height was still being
            # discovered let a two-line label in one cell be overwritten by the next row.
            lab_h = [est_height(it["label"], cw, SIZE["statlabel"], FONT_MONO, 1.35)
                     for it in items]
            nrow = -(-len(items) // ncol)
            row_h = [max([lab_h[i] for i in range(len(items)) if i // ncol == r] or [0])
                     + SIZE["stat"] * 1.05 + 8 for r in range(nrow)]
            # `for i, ...` here used to shadow the enclosing while-loop's cursor, so after a
            # statistics block the cursor sat past the end of the list and every later block
            # on that slide was silently dropped. Slide 9 lost its pull quote to it.
            for k, it in enumerate(items):
                col, row = k % ncol, k // ncol
                sx = x + col * (cw + GAP)
                sy = y + sum(row_h[:row]) + row * 20
                rule(slide, sx, sy, cw, RAIL["method"], 1.6)
                _, tf = box(slide, sx, sy + 8, cw, SIZE["stat"] * 1.15)
                put(tf.paragraphs[0], it["value"], size_px=SIZE["stat"], face=FONT_DISPLAY,
                    bold=True, line=1.02, strike=it["struck"],
                    color=MUTED if it["struck"] else INK)
                _, tf2 = box(slide, sx, sy + SIZE["stat"] * 1.05 + 8, cw, lab_h[k] + 4)
                put(tf2.paragraphs[0], it["label"], size_px=SIZE["statlabel"],
                    face=FONT_MONO, color=MUTED, line=1.35, spc=55, upper=True)
            y += sum(row_h) + nrow * 20 + 4

        elif b["t"] == "led":
            kw = max(est_height(i["key"], 9999, SIZE["key"], FONT_MONO, 1.0) and
                     len("".join(t for t, _ in i["key"])) for i in b["items"])
            keyw = min(w * 0.34, kw * SIZE["key"] * 0.60 + 10)
            for it in b["items"]:
                tw = w - keyw - 14
                h = max(est_height(it["text"], tw, SIZE["p"], FONT_BODY, LINE),
                        est_height(it["key"], keyw, SIZE["key"], FONT_MONO, 1.35))
                _, tf = box(slide, x, y + 2, keyw, h)
                put(tf.paragraphs[0], it["key"], size_px=SIZE["key"], face=FONT_MONO,
                    color=ACCENT, line=1.35, spc=80, upper=True)
                _, tf2 = box(slide, x + keyw + 14, y, tw, h + 4)
                put(tf2.paragraphs[0], it["text"], size_px=SIZE["p"], line=LINE)
                y += h + 12

        elif b["t"] == "table":
            y = _table(slide, b, x, y, w)

        elif b["t"] == "fig":
            y = _figure(slide, b, x, y, w, avail - (y - TOP))

        elif b["t"] == "h":
            size = SIZE["h1"] if b["level"] == 1 else SIZE["h2"]
            runs = _runs(b["html"])
            h = est_height(runs, w, size, FONT_DISPLAY, 1.02)
            _, tf = box(slide, x, y, w, h + 6)
            put(tf.paragraphs[0], runs, size_px=size, face=FONT_DISPLAY, bold=True,
                line=1.02)
            y += h + 12
    return y


def _table(slide, b, x, y, w) -> float:
    head, rows = b["head"], b["rows"]
    ncol = len(head)
    # Column widths: the deck lets the browser measure. Here they are set from the content,
    # because a PowerPoint table with equal columns puts a twelve-character hash in the same
    # space as a fifty-character study title.
    widths = []
    for i in range(ncol):
        longest = max([len("".join(t for t, _ in head[i]))]
                      + [len("".join(t for t, _ in r[i]["runs"])) for r in rows
                         if i < len(r)])
        widths.append(max(4, longest))
    total = sum(widths)
    widths = [w * v / total for v in widths]
    rh = SIZE["table"] * 1.9
    shape = slide.shapes.add_table(len(rows) + 1, ncol, px(x), px(y), px(w),
                                   px(rh * (len(rows) + 1)))
    tbl = shape.table
    tbl.first_row = False
    for i, cw in enumerate(widths):
        tbl.columns[i].width = px(cw)
    for r in range(len(rows) + 1):
        tbl.rows[r].height = px(rh)
    # python-pptx has no API for "no banding, no theme fill", so strip the style reference.
    tbl._tbl.tblPr.set("bandRow", "0")
    tbl._tbl.tblPr.set("firstRow", "0")

    def fill(cell, runs, *, header=False, mono=False, align=PP_ALIGN.LEFT):
        cell.margin_left = cell.margin_right = px(5)
        cell.margin_top = cell.margin_bottom = px(2)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL
        put(cell.text_frame.paragraphs[0], runs,
            size_px=SIZE["table"] * (0.82 if header else 1.0),
            face=FONT_MONO if (mono or header) else FONT_BODY,
            color=MUTED if header else INK, line=1.15, spc=80 if header else None,
            upper=header, align=align)

    for i, h in enumerate(head):
        fill(tbl.cell(0, i), h, header=True)
    for r, row in enumerate(rows, 1):
        for i, cell in enumerate(row):
            cls = cell["attrs"].get("class", "")
            style = cell["attrs"].get("style", "")
            fill(tbl.cell(r, i), cell["runs"], mono="h" in cls.split(),
                 align=PP_ALIGN.RIGHT if "right" in style else PP_ALIGN.LEFT)
    return y + rh * (len(rows) + 1) + 14


def _figure(slide, b, x, y, w, avail) -> float:
    if not b["png"]:
        return y
    from PIL import Image
    img = Image.open(io.BytesIO(b["png"]))
    iw, ih = img.size
    cap_h = est_height(b["caption"], w, SIZE["cap"], FONT_MONO, 1.45) if b["caption"] else 0
    room = max(60.0, avail - cap_h - 16)
    scale = min(w / iw, room / ih)
    dw, dh = iw * scale, ih * scale
    slide.shapes.add_picture(io.BytesIO(b["png"]), px(x + (w - dw) / 2), px(y),
                             px(dw), px(dh))
    y += dh + 10
    if cap_h:
        _, tf = box(slide, x, y, w, cap_h + 4)
        put(tf.paragraphs[0], b["caption"], size_px=SIZE["cap"], face=FONT_MONO,
            color=MUTED, line=1.45)
        y += cap_h + 6
    return y


def _title_slide(slide, blocks) -> None:
    """The opening and closing slides: eyebrow, headline, thesis, identifiers, footer."""
    y = _eyebrow(slide, next(b["html"] for b in blocks if b["t"] == "eyebrow"), TOP)
    h = next(b for b in blocks if b["t"] == "h")
    runs = _runs(h["html"])
    size = SIZE["h1"] if len(" ".join(t for t, _ in runs)) < 30 else SIZE["h1"] * 0.72
    hh = est_height(runs, CW * 0.82, size, FONT_DISPLAY, 1.02)
    _, tf = box(slide, L, y + 34, CW * 0.82, hh + 8)
    put(tf.paragraphs[0], runs, size_px=size, face=FONT_DISPLAY, bold=True, line=1.02)
    y2 = y + 34 + hh + 18

    thesis = next((b for b in blocks if b["t"] == "p"), None)
    if thesis:
        th = est_height(thesis["runs"], CW * 0.46, SIZE["thesis"], FONT_BODY, 1.42)
        _, tf2 = box(slide, L, y2, CW * 0.46, th + 6)
        put(tf2.paragraphs[0], thesis["runs"], size_px=SIZE["thesis"], line=1.42)

    meta = next((b for b in blocks if b["t"] == "meta"), None)
    if meta:
        def _line_runs(fragment: str):
            spans = re.findall(r"<span[^>]*>(.*?)</span>", fragment, re.S)
            if len(spans) > 1:                       # a flex row of identifiers
                out = []
                for i2, sp in enumerate(spans):
                    if i2:
                        out.append(("   ·   ", frozenset()))
                    out.extend(_runs(sp))
                return out
            return _runs(fragment)

        lines = [_line_runs(m)
                 for m in re.findall(r"<div[^>]*>(.*?)</div>", meta["html"], re.S)]
        if not lines:
            lines = [_runs(meta["html"])]
        mh = sum(est_height(r, CW, SIZE["meta"], FONT_MONO, 1.9) for r in lines)
        _, tf3 = box(slide, L, FOOT_Y - 22 - mh, CW, mh + 6)
        for i, r in enumerate(lines):
            p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
            put(p, r, size_px=SIZE["meta"], face=FONT_MONO, color=MUTED, line=1.9, spc=60)
    _footer(slide, next(b["html"] for b in blocks if b["t"] == "footer"))


def _visible_text(html: str) -> str:
    """Everything a reader sees in a slide's markup, with tags and attributes removed."""
    html = re.sub(r"<img[^>]*>", " ", html)
    html = re.sub(r"<[^>]+>", " ", html)
    return re.sub(r"\s+", " ", _Inline.unescape(html))


def _assert_nothing_dropped(prs, html: str) -> None:
    """Every word the deck shows must appear on the matching PowerPoint slide.

    A parser that silently skips a block produces a file that looks finished and is not.
    That happened once already: a shadowed loop cursor made every block after a statistics
    row disappear, and slide 9 shipped without its conclusion. Geometry checks would not
    have caught it, because the missing shape has no geometry.
    """
    sections = [_tag_span(html, m.start())[0]
                for m in re.finditer(r'<section class="slide', html)]
    missing: list[str] = []
    for n, (sec, slide) in enumerate(zip(sections, prs.slides), 1):
        # Case-insensitive: the deck sets eyebrows, statistic labels, definition keys and
        # table headers in small caps via CSS, and this generator applies that with
        # upper=True, so "Premise" legitimately becomes "PREMISE" on the slide.
        want = set(re.findall(r"[a-z][a-z'-]{4,}", _visible_text(sec).lower()))
        got = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        got += " ".join(c.text for sh in slide.shapes if sh.has_table
                        for r in sh.table.rows for c in r.cells)
        have = set(re.findall(r"[a-z][a-z'-]{4,}", got.lower()))
        gone = sorted(want - have)
        if gone:
            missing.append(f"slide {n}: {len(gone)} word(s) not carried over, "
                           f"e.g. {gone[:8]}")
    if missing:
        raise AssertionError(
            "the .pptx does not contain everything the deck shows:\n  "
            + "\n  ".join(missing))


def build(src: Path | None = None, out: Path | None = None) -> int:
    global FONT_DISPLAY, FONT_BODY, FONT_MONO
    # A relative path on the command line is relative to the caller's directory, but every
    # message this module prints is relative to the repository root. Resolve once, here.
    src = (Path(src).resolve() if src else DEFAULT_SRC)
    out = (Path(out).resolve() if out else src.with_suffix(".pptx"))
    if not src.exists():
        raise FileNotFoundError(f"{src} is missing; build the HTML deck first")
    raw = src.read_text()
    if 'data-lang="ko"' in raw:
        FONT_DISPLAY, FONT_BODY, FONT_MONO = FONT_KO_DISPLAY, FONT_KO_BODY, FONT_KO_MONO
    slides = parse_slides(raw)

    prs = Presentation()
    prs.slide_width, prs.slide_height = px(DECK_W), px(DECK_H)
    blank = prs.slide_layouts[6]

    for s in slides:
        sl = prs.slides.add_slide(blank)
        bg = sl.background.fill
        bg.solid()
        bg.fore_color.rgb = PANEL
        rail = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px(5), px(DECK_H))
        rail.fill.solid()
        rail.fill.fore_color.rgb = RAIL.get(s["kind"], MUTED)
        rail.line.fill.background()
        rail.shadow.inherit = False

        if s["title_slide"]:
            _title_slide(sl, s["blocks"])
            continue

        y = _header(sl, s["blocks"], TOP)
        body = [b for b in s["blocks"] if b["t"] not in ("header", "footer")]
        _blocks(sl, body, L, y, CW, FOOT_Y - 18 - y + (y - TOP))
        _footer(sl, next(b["html"] for b in s["blocks"] if b["t"] == "footer"))

    prs.core_properties.title = "CognitionBioChem"
    prs.core_properties.author = "Seung Ho Jung"
    # The core-properties fields cap at 255 characters, so this says where the file came
    # from and stops; the reasoning lives in this module's docstring.
    prs.core_properties.comments = (
        f"Generated by platform/build_pptx.py from {src.name} at {git_sha()}. Every number "
        f"on these slides is read from a study artefact at build time -- re-run the "
        f"generators rather than retyping one.")
    _assert_nothing_dropped(prs, raw)
    prs.save(out)
    shown = out.relative_to(REPO) if out.is_relative_to(REPO) else out
    print(f"wrote {shown}")
    print(f"  slides: {len(slides)}   size: {out.stat().st_size / 1024:.0f} KB")
    print(f"  fonts:  {FONT_DISPLAY} / {FONT_BODY} / {FONT_MONO}")
    return 0


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("-")]
    raise SystemExit(build(*(Path(a) for a in args[:2])))
